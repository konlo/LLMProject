# app_ppt_edit.py
import os
import io
import json
from typing import List, Dict, Any

import streamlit as st
from dotenv import load_dotenv

# LLM (Gemini via LangChain)
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.schema import SystemMessage, HumanMessage

# PPTX editing
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# -----------------------------
# 1) 기본 설정
# -----------------------------
load_dotenv()

st.set_page_config(page_title="PPT Editor (Gemini)", page_icon="🖼️")
st.title("🖼️ PPT Editor (Gemini + python-pptx)")
st.caption("Upload a .pptx, describe your changes in natural language, and apply them.")

GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    st.error("환경변수 GOOGLE_API_KEY 가 설정되어 있지 않습니다. .env 또는 OS 환경변수로 설정하세요.")
    st.stop()

llm = ChatGoogleGenerativeAI(
    model="gemini-2.5-flash",
    api_key=GOOGLE_API_KEY,
    temperature=0.2,
    streaming=True,
)

# -----------------------------
# 2) 파일 업로더 + 프롬프트
# -----------------------------
uploaded_ppt = st.file_uploader("① PPTX 파일을 업로드하세요", type=["pptx"])
uploaded_images = st.file_uploader("선택) 이미지 파일 업로드 (여러 개 가능)", type=["png", "jpg", "jpeg"], accept_multiple_files=True)

prompt = st.text_area(
    "② 프롬프트 (예: '모든 슬라이드에서 회사명을 삼성전자→Samsung으로 변경하고, 마지막에 요약 슬라이드 추가, 1번 슬라이드 제목을 “프로젝트 개요”로')",
    height=130,
)

col_run, col_show = st.columns([1,1])
run = col_run.button("Run ▶️")
show_plan_only = col_show.checkbox("수정 전 LLM 작업계획(JSON)만 먼저 보고 싶어요", value=False)

# -----------------------------
# 3) LLM에게 “작업계획 JSON” 생성 지시
#    - 모델은 아래 스키마로만 JSON을 출력
# -----------------------------
SYSTEM_INSTR = """You are a PPT editing planner. Convert the user's request into a STRICT JSON plan.
Output MUST be valid JSON ONLY, no prose, no code fences.

JSON schema:
{
  "actions": [
    {
      "type": "replace_text",
      "params": {
        "find": "string",           // 텍스트 검색 (부분 포함)
        "replace": "string",        // 대체 텍스트
        "match_case": false         // 대소문자 구분 여부
      }
    },
    {
      "type": "set_title",
      "params": {
        "slide_index": 0,           // 0-based index
        "text": "string",
        "align": "left|center|right"// optional
      }
    },
    {
      "type": "add_title_content_slide",
      "params": {
        "title": "string",
        "bullets": ["string", "string", "..."]  // optional
      }
    },
    {
      "type": "add_image",
      "params": {
        "slide_index": 1,          // 어느 슬라이드에 넣을지
        "image_name": "string",    // 업로드한 이미지 파일명 (정확히 일치)
        "left_in": 1.0,            // inches
        "top_in": 1.5,
        "width_in": 4.0            // optional (없으면 원본 비율)
      }
    },
    {
      "type": "append_summary_slide",
      "params": {
        "title": "Summary",                // optional
        "points": ["string", "string"]     // optional
      }
    }
  ]
}
If the user's request is unclear, make reasonable assumptions and prefer generic actions above.
Do not invent new action types. Only use the listed types.
"""

def ask_llm_for_plan(user_prompt: str) -> Dict[str, Any]:
    messages = [
        SystemMessage(content=SYSTEM_INSTR),
        HumanMessage(content=user_prompt),
    ]
    resp = llm.invoke(messages)
    # resp.content should be JSON
    content = resp.content
    # 방어적으로 JSON만 추출
    try:
        # content 전체가 JSON이라고 가정
        data = json.loads(content)
        return data
    except Exception:
        # 괄호 범위 추출 시도
        try:
            start = content.find("{")
            end = content.rfind("}")
            if start != -1 and end != -1:
                data = json.loads(content[start:end+1])
                return data
        except Exception:
            pass
    raise ValueError("LLM이 유효한 JSON을 생성하지 못했습니다.\n\n응답:\n" + content)


# -----------------------------
# 4) PPTX 편집 유틸
# -----------------------------
def iter_text_shapes(slide):
    for shape in slide.shapes:
        if not hasattr(shape, "has_text_frame"):
            continue
        if not shape.has_text_frame:
            continue
        yield shape

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}

def apply_replace_text(prs: Presentation, params: Dict[str, Any], logs: List[str]):
    find = params.get("find", "")
    replace = params.get("replace", "")
    match_case = params.get("match_case", False)

    if not find:
        logs.append("replace_text: 'find' 누락 → skip")
        return

    count = 0
    for slide in prs.slides:
        for shape in iter_text_shapes(slide):
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    src = run.text
                    if not match_case:
                        if find.lower() in src.lower():
                            # 대소문자 무시 치환
                            # 간단치환: 전체 일치가 아니라 포함이므로, 보존을 위해 replace를 case-insensitive하게 처리
                            # 여기서는 단순화: 소문자 비교 후 위치 매칭 어려우니 파라 전체를 replace
                            # 정확 치환 원하면 정규식/대소문자 보존 로직 추가
                            run.text = src.replace(find, replace).replace(find.capitalize(), replace).replace(find.lower(), replace)
                            count += 1
                    else:
                        if find in src:
                            run.text = src.replace(find, replace)
                            count += 1
    logs.append(f"replace_text: '{find}' → '{replace}' ({count} runs)")

def apply_set_title(prs: Presentation, params: Dict[str, Any], logs: List[str]):
    idx = int(params.get("slide_index", 0))
    text = params.get("text", "")
    align = params.get("align")

    if idx < 0 or idx >= len(prs.slides):
        logs.append(f"set_title: slide_index {idx} 범위 밖 → skip")
        return

    slide = prs.slides[idx]

    # 타이틀 플레이스홀더 우선
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape:
        title_shape.text = text
        if align and align in ALIGN_MAP:
            for p in title_shape.text_frame.paragraphs:
                p.alignment = ALIGN_MAP[align]
        logs.append(f"set_title: slide {idx} 제목을 '{text}'로 설정")
        return

    # 없으면 첫 텍스트 shape 사용
    for shape in iter_text_shapes(slide):
        shape.text = text
        if align and align in ALIGN_MAP:
            for p in shape.text_frame.paragraphs:
                p.alignment = ALIGN_MAP[align]
        logs.append(f"set_title: slide {idx} 첫 텍스트 상자에 제목 대입 ('{text}')")
        return

    logs.append(f"set_title: slide {idx} 텍스트 상자 없음 → skip")

def apply_add_title_content_slide(prs: Presentation, params: Dict[str, Any], logs: List[str]):
    title = params.get("title", "New Slide")
    bullets = params.get("bullets", [])

    # 대부분 템플릿의 레이아웃 1이 Title and Content
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # 타이틀
    if slide.shapes.title:
        slide.shapes.title.text = title
    # 본문 placeholder
    ph = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 2:  # BODY
            ph = shape
            break
    if ph and ph.has_text_frame:
        tf = ph.text_frame
        tf.clear()
        first = True
        for b in bullets:
            if first:
                tf.text = b
                first = False
            else:
                p = tf.add_paragraph()
                p.text = b
    logs.append(f"add_title_content_slide: '{title}' (bullets: {len(bullets)})")

def apply_add_image(prs: Presentation, params: Dict[str, Any], logs: List[str], image_map: Dict[str, bytes]):
    idx = int(params.get("slide_index", len(prs.slides)-1))
    image_name = params.get("image_name")
    left_in = float(params.get("left_in", 1.0))
    top_in = float(params.get("top_in", 1.0))
    width_in = params.get("width_in", None)

    if idx < 0 or idx >= len(prs.slides):
        logs.append(f"add_image: slide_index {idx} 범위 밖 → skip")
        return
    if not image_name or image_name not in image_map:
        logs.append(f"add_image: '{image_name}'를 찾을 수 없음 → skip (업로드 이미지 파일명 확인)")
        return

    slide = prs.slides[idx]
    image_bytes = image_map[image_name]
    stream = io.BytesIO(image_bytes)

    left = Inches(left_in)
    top = Inches(top_in)
    if width_in:
        slide.shapes.add_picture(stream, left, top, width=Inches(float(width_in)))
    else:
        slide.shapes.add_picture(stream, left, top)
    logs.append(f"add_image: slide {idx}에 '{image_name}' 삽입")

def apply_append_summary_slide(prs: Presentation, params: Dict[str, Any], logs: List[str]):
    title = params.get("title", "Summary")
    points = params.get("points", [])

    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        slide.shapes.title.text = title

    # 본문 placeholder
    ph = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 2:  # BODY
            ph = shape
            break
    if ph and ph.has_text_frame:
        tf = ph.text_frame
        tf.clear()
        if points:
            tf.text = points[0]
            for ptxt in points[1:]:
                p = tf.add_paragraph()
                p.text = ptxt
    logs.append(f"append_summary_slide: '{title}' (points: {len(points)})")

ACTION_TABLE = {
    "replace_text": apply_replace_text,
    "set_title": apply_set_title,
    "add_title_content_slide": apply_add_title_content_slide,
    "add_image": apply_add_image,
    "append_summary_slide": apply_append_summary_slide,
}

# -----------------------------
# 5) 실행
# -----------------------------
if run:
    if not uploaded_ppt:
        st.warning("PPTX 파일을 먼저 업로드하세요.")
        st.stop()
    if not prompt.strip():
        st.warning("프롬프트를 입력하세요.")
        st.stop()

    # 5-1) LLM 계획 수립
    with st.spinner("LLM이 작업 계획(JSON)을 생성 중..."):
        plan = ask_llm_for_plan(prompt)

    st.subheader("LLM 작업 계획 (JSON)")
    st.code(json.dumps(plan, ensure_ascii=False, indent=2), language="json")
    if show_plan_only:
        st.info("체크박스를 껴서, 우선 작업 계획만 확인하도록 설정했습니다. 체크 해제 후 다시 Run 하세요.")
        st.stop()

    # 5-2) 이미지 맵 구성 (이름 → bytes)
    image_map = {}
    if uploaded_images:
        for img in uploaded_images:
            image_map[img.name] = img.read()

    # 5-3) PPT 적용
    prs = Presentation(uploaded_ppt)
    logs: List[str] = []

    actions = plan.get("actions", [])
    if not isinstance(actions, list):
        st.error("LLM JSON에 'actions' 리스트가 없습니다.")
        st.stop()

    for i, action in enumerate(actions, start=1):
        atype = action.get("type")
        params = action.get("params", {})
        fn = ACTION_TABLE.get(atype)
        if not fn:
            logs.append(f"[{i}] 지원하지 않는 type: {atype} → skip")
            continue

        try:
            if atype == "add_image":
                fn(prs, params, logs, image_map)  # 이미지만 별도 시그니처
            else:
                fn(prs, params, logs)
        except Exception as e:
            logs.append(f"[{i}] {atype} 적용 중 오류: {e}")

    # 5-4) 결과 저장 & 다운로드
    out_buffer = io.BytesIO()
    prs.save(out_buffer)
    out_buffer.seek(0)

    st.subheader("적용 로그")
    st.write("\n".join(logs) if logs else "적용된 로그가 없습니다.")

    st.download_button(
        label="수정된 PPTX 다운로드",
        data=out_buffer,
        file_name="edited.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

# -----------------------------
# 6) 하단 도움말
# -----------------------------
with st.expander("프롬프트 예시 보기"):
    st.markdown("""
- **전역 치환**: "모든 슬라이드에서 '삼성전자'를 'Samsung'으로 바꿔줘"
- **제목 변경**: "1번 슬라이드 제목을 '프로젝트 개요'로, 가운데 정렬로 바꿔"
- **슬라이드 추가**: "마지막에 요약 슬라이드를 추가하고 항목은 3개로 만들어"
- **타이틀+본문 슬라이드**: "새 슬라이드를 추가해서 제목 '로드맵', 본문 불릿 3개(1분기, 2분기, 3분기)"
- **이미지 삽입**: "2번 슬라이드에 업로드한 'logo.png'를 왼쪽 1.0인치, 위 1.0인치, 폭 2.0인치로 넣어"
""")
