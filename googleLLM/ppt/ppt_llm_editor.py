#!/usr/bin/env python3
"""
ppt_llm_editor.py
-----------------
Edit a local PowerPoint (.pptx) using instructions proposed by an LLM, but **all edits are applied locally**
via python-pptx (the LLM never touches the binary directly).

Workflow
1) Extract a structured summary of slides (texts, tables) from the source PPTX.
2) Send the summary + your high-level instructions to an LLM and ask it to output a JSON "edit plan"
   matching the schema below.
3) Apply the JSON edits to the PPTX locally and save a new file.

Supported operations in the edit plan (JSON):
- replace_text
- add_bullets
- edit_table_cell
- delete_slide
- add_image
- add_new_slide (title + bullet content)

Limitations
- Moving/reordering slides is not supported by python-pptx API directly (workarounds are clunky).
- "Global" style/theme changes are limited; we can set font attributes for individual shapes.
- Image replacement requires a local file path to the new image.
"""

from __future__ import annotations
import json
from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Union

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor


# -----------------------------
# 1) Utilities: extraction
# -----------------------------

def extract_presentation_structure(pptx_path: str) -> Dict[str, Any]:
    """
    Return a lightweight JSON-serializable structure of the deck:
    slides -> shapes (text frames) and tables (cell text grid).
    """
    prs = Presentation(pptx_path)
    summary = {"num_slides": len(prs.slides), "slides": []}

    for s_idx, slide in enumerate(prs.slides):
        slide_info = {"index": s_idx, "shapes": [], "tables": []}

        for sh_idx, shape in enumerate(slide.shapes):
            # Text shapes
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                texts = []
                for p in shape.text_frame.paragraphs:
                    run_text = "".join(run.text for run in p.runs) or p.text
                    texts.append(run_text)
                slide_info["shapes"].append({
                    "shape_idx": sh_idx,
                    "type": "text",
                    "text": "\n".join(texts).strip()
                })

            # Tables
            if shape.has_table:
                tbl = shape.table
                table_grid = []
                for r in range(len(tbl.rows)):
                    row_vals = []
                    for c in range(len(tbl.columns)):
                        row_vals.append(tbl.cell(r, c).text)
                    table_grid.append(row_vals)
                slide_info["tables"].append({
                    "table_idx": sh_idx,
                    "rows": len(tbl.rows),
                    "cols": len(tbl.columns),
                    "data": table_grid,
                })

        summary["slides"].append(slide_info)

    return summary


# -----------------------------
# 2) Edit Plan Schema
# -----------------------------

@dataclass
class ReplaceTextOp:
    op: str  # "replace_text"
    slide_index: int
    find: str
    replace: str
    shape_idx: Optional[int] = None  # if None, search all text shapes on the slide
    case_sensitive: bool = False
    first_only: bool = False


@dataclass
class AddBulletsOp:
    op: str  # "add_bullets"
    slide_index: int
    shape_idx: int
    bullets: List[str]
    mode: str = "append"  # "append" or "replace"
    level: int = 0  # bullet level


@dataclass
class EditTableCellOp:
    op: str  # "edit_table_cell"
    slide_index: int
    table_idx: int
    row: int
    col: int
    new_text: str


@dataclass
class DeleteSlideOp:
    op: str  # "delete_slide"
    slide_index: int


@dataclass
class AddImageOp:
    op: str  # "add_image"
    slide_index: int
    image_path: str
    left_in: float = 1.0
    top_in: float = 1.5
    width_in: Optional[float] = 4.0  # set width; height auto-scales if None
    height_in: Optional[float] = None


@dataclass
class AddNewSlideOp:
    op: str  # "add_new_slide"
    layout: str  # "title_and_content" | "title_only" | etc.
    title: str
    bullets: Optional[List[str]] = None
    align_center: bool = False


EditOp = Union[ReplaceTextOp, AddBulletsOp, EditTableCellOp, DeleteSlideOp, AddImageOp, AddNewSlideOp]


# -----------------------------
# 3) Apply edits locally
# -----------------------------

def _delete_slide(prs: Presentation, index: int) -> None:
    """
    Delete slide by index using private API (works in practice).
    """
    slide_id_list = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slides = list(slide_id_list)
    slide_id_list.remove(slides[index])


def apply_edits_to_pptx(pptx_path: str, edits: List[Dict[str, Any]], out_path: str) -> None:
    prs = Presentation(pptx_path)

    for i, op in enumerate(edits):
        kind = op.get("op")
        if kind == "replace_text":
            slide_index = op["slide_index"]
            find = op["find"]
            replace = op["replace"]
            shape_idx = op.get("shape_idx")
            case_sensitive = op.get("case_sensitive", False)
            first_only = op.get("first_only", False)

            slide = prs.slides[slide_index]
            replaced_once = False

            def do_replace(haystack: str, needle: str, replacement: str, cs: bool) -> str:
                if cs:
                    return haystack.replace(needle, replacement, 1 if first_only else -1)
                # case-insensitive
                import re
                flags = re.IGNORECASE
                return re.sub(re.escape(needle), replacement, haystack, count=1 if first_only else 0, flags=flags)

            for idx, shape in enumerate(slide.shapes):
                if shape_idx is not None and idx != shape_idx:
                    continue
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    tf = shape.text_frame
                    # concatenate paragraphs to a long string, replace, then write back as one paragraph
                    original_text = "\\n".join(p.text for p in tf.paragraphs)
                    new_text = do_replace(original_text, find, replace, case_sensitive)
                    if new_text != original_text:
                        # clear and write back
                        for _ in range(len(tf.paragraphs) - 1):
                            p = tf.paragraphs[-1]
                            p._element.getparent().remove(p._element)  # remove extra paragraphs
                        tf.paragraphs[0].text = new_text
                        if first_only:
                            replaced_once = True
                if first_only and replaced_once:
                    break

        elif kind == "add_bullets":
            slide_index = op["slide_index"]
            shape_idx = op["shape_idx"]
            bullets = op.get("bullets", [])
            mode = op.get("mode", "append")
            level = op.get("level", 0)

            slide = prs.slides[slide_index]
            shape = slide.shapes[shape_idx]
            if not (hasattr(shape, "has_text_frame") and shape.has_text_frame):
                continue
            tf = shape.text_frame

            if mode == "replace":
                # clear all paragraphs
                for _ in range(len(tf.paragraphs) - 1):
                    p = tf.paragraphs[-1]
                    p._element.getparent().remove(p._element)
                tf.paragraphs[0].text = ""
                tf.paragraphs[0].level = level

            # append bullets
            for b in bullets:
                p = tf.add_paragraph() if tf.paragraphs[0].text != "" else tf.paragraphs[0]
                if p.text != "":
                    p = tf.add_paragraph()
                p.text = b
                p.level = level

        elif kind == "edit_table_cell":
            slide_index = op["slide_index"]
            table_idx = op["table_idx"]
            row = op["row"]
            col = op["col"]
            new_text = op["new_text"]
            slide = prs.slides[slide_index]
            shape = slide.shapes[table_idx]
            if not shape.has_table:
                continue
            tbl = shape.table
            tbl.cell(row, col).text = new_text

        elif kind == "delete_slide":
            slide_index = op["slide_index"]
            _delete_slide(prs, slide_index)

        elif kind == "add_image":
            slide_index = op["slide_index"]
            image_path = op["image_path"]
            left_in = float(op.get("left_in", 1.0))
            top_in = float(op.get("top_in", 1.5))
            width_in = op.get("width_in")
            height_in = op.get("height_in")
            slide = prs.slides[slide_index]
            left = Inches(left_in)
            top = Inches(top_in)
            if width_in is not None:
                width = Inches(float(width_in))
                slide.shapes.add_picture(image_path, left, top, width=width)
            elif height_in is not None:
                height = Inches(float(height_in))
                slide.shapes.add_picture(image_path, left, top, height=height)
            else:
                slide.shapes.add_picture(image_path, left, top)

        elif kind == "add_new_slide":
            layout_name = op.get("layout", "title_and_content")
            title = op.get("title", "")
            bullets = op.get("bullets", [])
            align_center = op.get("align_center", False)

            # Map a few common layout names to indices (depends on template; best-effort)
            layout_map = {
                "title_and_content": 1,
                "title_only": 5,
                "blank": 6,
            }
            layout_idx = layout_map.get(layout_name, 1)
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)

            # Title
            if slide.shapes.title:
                slide.shapes.title.text = title

            # Content placeholder (if exists)
            for shape in slide.shapes:
                if shape.is_placeholder and hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    tf = shape.text_frame
                    tf.clear()
                    if bullets:
                        first = True
                        for b in bullets:
                            p = tf.paragraphs[0] if first else tf.add_paragraph()
                            p.text = b
                            if align_center:
                                p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                            first = False
                    break  # only fill the first content placeholder

        else:
            print(f"[WARN] Unknown op: {kind}; skipping")

    prs.save(out_path)
    print(f"[OK] Saved: {out_path}")


# -----------------------------
# 4) LLM interface (stub)
# -----------------------------

LLM_PROMPT_TEMPLATE = """\
You are a presentation editor. Given a JSON description of a PPTX deck and the user's goals,
produce a STRICT JSON edit plan following the provided schema.
Only output valid JSON with a top-level key "edits": [ ... ].

User goals:
{user_goals}

Deck summary (JSON):
{deck_summary}

JSON schema (operations you can emit):
- replace_text: {{"op":"replace_text","slide_index":int,"find":str,"replace":str,"shape_idx":int|null,"case_sensitive":false,"first_only":false}}
- add_bullets: {{"op":"add_bullets","slide_index":int,"shape_idx":int,"bullets":[str,...],"mode":"append|replace","level":0}}
- edit_table_cell: {{"op":"edit_table_cell","slide_index":int,"table_idx":int,"row":int,"col":int,"new_text":str}}
- delete_slide: {{"op":"delete_slide","slide_index":int}}
- add_image: {{"op":"add_image","slide_index":int,"image_path":str,"left_in":float,"top_in":float,"width_in":float|null,"height_in":float|null}}
- add_new_slide: {{"op":"add_new_slide","layout":"title_and_content","title":str,"bullets":[str,...]}}

Respond ONLY with JSON like:
{{"edits":[{{...}},{{...}}]}}
"""

def build_llm_prompt(deck_summary: Dict[str, Any], user_goals: str) -> str:
    return LLM_PROMPT_TEMPLATE.format(
        user_goals=user_goals,
        deck_summary=json.dumps(deck_summary, ensure_ascii=False, indent=2)
    )


def call_llm_and_get_edits(prompt: str) -> Dict[str, Any]:
    """
    Implement this to call your preferred LLM.
    Example for OpenAI:
        from openai import OpenAI
        client = OpenAI()
        resp = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role":"user","content": prompt}],
            response_format={"type": "json_object"},
            temperature=0
        )
        content = resp.choices[0].message.content
        return json.loads(content)
    """
    raise NotImplementedError("Hook up your LLM provider here and return a dict with {'edits': [...]}")


# -----------------------------
# 5) CLI
# -----------------------------

def main():
    import argparse
    p = argparse.ArgumentParser(description="Edit a PPTX using an LLM-generated JSON edit plan (applied locally).")
    p.add_argument("--in", dest="in_path", required=True, help="Path to source .pptx")
    p.add_argument("--out", dest="out_path", required=True, help="Path to save edited .pptx")
    p.add_argument("--goals", dest="goals", required=False, default="Improve clarity; fix typos; add one summary slide.",
                   help="High-level instructions for the LLM")
    p.add_argument("--dump-structure", action="store_true", help="Only extract deck structure to JSON and exit")
    p.add_argument("--apply-edits-json", dest="edits_json", help="Apply edits from a local JSON file (skip LLM call)")
    args = p.parse_args()

    # 1) Extract structure
    deck_summary = extract_presentation_structure(args.in_path)

    if args.dump_structure:
        print(json.dumps(deck_summary, ensure_ascii=False, indent=2))
        return

    # 2) Either load edits from file, or ask an LLM (user must implement call_llm_and_get_edits)
    if args.edits_json:
        with open(args.edits_json, "r", encoding="utf-8") as f:
            edit_plan = json.load(f)
    else:
        prompt = build_llm_prompt(deck_summary, args.goals)
        edit_plan = call_llm_and_get_edits(prompt)  # <-- implement

    # 3) Apply
    edits = edit_plan.get("edits", [])
    if not isinstance(edits, list):
        raise ValueError("Edit plan must contain a list at key 'edits'.")
    apply_edits_to_pptx(args.in_path, edits, args.out_path)


if __name__ == "__main__":
    main()
